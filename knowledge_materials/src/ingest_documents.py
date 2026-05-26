"""営業資料 PPTX を documents / chunks コンテナに取り込むナレッジベース構築スクリプト.

- documents: 資料メタ情報（PK = /type）
- chunks   : 本文チャンク + 1536 次元埋め込み（PK = /document_id, vector policy 付き）

アプリ（agent-first-meeting）の CaseSearchPlugin がこの documents/chunks を読むため、
接続情報・埋め込み次元はアプリと統一する（agent_first_meeting.config.settings 経由、
DB は settings.cosmos_database = sales-agent、1536 次元）。chunks は VectorDistance
検索にベクトルポリシーが必須なので、vector-policy.json / index-policy.json を使って
「無ければ作成」する。data/ 配下のサンプル PPTX 3 点を投入する。

冪等性: doc_id はファイルパスから決定的に導出し、毎回 documents / chunks の
既存アイテムを全削除してから取り込む。これにより何度再実行しても重複が増えず、
スキーマ変更（slide_number の付け方など）も常に最新の状態へ入れ替わる。
"""
import hashlib
import json
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path

from azure.cosmos import CosmosClient, PartitionKey
from openai import AzureOpenAI
from pptx import Presentation

from agent_first_meeting.config import settings

sys.stdout.reconfigure(encoding="utf-8")

_ROOT = Path(__file__).resolve().parent.parent  # knowledge_materials/
_DATA_DIR = _ROOT / "data"
_VECTOR_POLICY = json.loads((_ROOT / "vector-policy.json").read_text("utf-8"))
_INDEX_POLICY = json.loads((_ROOT / "index-policy.json").read_text("utf-8"))

_CHUNK_SIZE = 500
_CHUNK_OVERLAP = 100

SAMPLE_FILES = [
    {
        "path": "proposal_manufacturing_2024.pptx",
        "title": "製造業向け生産管理DX提案書 2024年版",
        "type": "proposal",
        "industry": "製造業",
        "tags": ["DX", "IoT", "コスト削減", "予防保全"],
    },
    {
        "path": "proposal_it_security_2024.pptx",
        "title": "中堅企業向けセキュリティ強化提案書 2024年版",
        "type": "proposal",
        "industry": "IT",
        "tags": ["セキュリティ", "MFA", "EDR", "ランサムウェア対策"],
    },
    {
        "path": "catalog_cloud_services_2024.pptx",
        "title": "クラウドサービス製品カタログ 2024",
        "type": "catalog",
        "industry": "全業種",
        "tags": ["クラウド", "IoT", "セキュリティ", "データ分析"],
    },
]


# ----------------------------------------
# ① PPTX からスライドごとに (実スライド番号, 本文) を抽出
#    スライド境界を保持し、後段で chunk に実スライド番号を付与できるようにする。
# ----------------------------------------
def extract_slides_from_pptx(file_path: Path) -> list[tuple[int, str]]:
    prs = Presentation(str(file_path))
    slides: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    parts.append(line)
        if parts:
            slides.append((i, "\n".join(parts)))
    return slides


# ----------------------------------------
# ② テキストをチャンクに分割（500 文字 / 100 文字オーバーラップ・依存なし）
# ----------------------------------------
def split_into_chunks(text: str) -> list[str]:
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + _CHUNK_SIZE, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - _CHUNK_OVERLAP
    return chunks


# ----------------------------------------
# ③ OpenAI でベクトル生成（アプリと同じ 1536 次元）
# ----------------------------------------
def get_embedding(openai_client: AzureOpenAI, text: str) -> list[float]:
    return (
        openai_client.embeddings.create(
            model=settings.azure_openai_embedding_deployment,
            input=text,
            dimensions=1536,
        )
        .data[0]
        .embedding
    )


def _doc_id_for(path: str) -> str:
    """ファイルパスから決定的な doc_id を導出する（再取り込みで重複させないため）."""
    digest = hashlib.sha1(path.encode("utf-8")).hexdigest()[:16]
    return f"doc_{digest}"


def _purge_all(container, pk_field: str, label: str) -> int:
    """コンテナ内の全アイテムを削除する（クリーンな再取り込みのため）.

    pk_field はそのコンテナのパーティションキーに対応するフィールド名
    （documents は 'type'、chunks は 'document_id'）。削除に必要な id と
    パーティションキーだけを射影し、chunks の埋め込みベクトルまで取得しない。
    """
    items = list(
        container.query_items(
            query=f"SELECT c.id, c.{pk_field} AS pk FROM c",
            enable_cross_partition_query=True,
        )
    )
    for item in items:
        container.delete_item(item=item["id"], partition_key=item["pk"])
    if items:
        print(f"  ↻ {label}: 既存 {len(items)} 件を削除（クリーン再取り込み）")
    return len(items)


def _setup_containers(db):
    """documents (/type) と chunks (/document_id + vector policy) を無ければ作成する."""
    documents = db.create_container_if_not_exists(
        id="documents",
        partition_key=PartitionKey(path="/type"),
    )
    chunks = db.create_container_if_not_exists(
        id="chunks",
        partition_key=PartitionKey(path="/document_id"),
        indexing_policy=_INDEX_POLICY,
        vector_embedding_policy=_VECTOR_POLICY,
    )
    return documents, chunks


# ----------------------------------------
# メイン：1 ファイルを投入する
# ----------------------------------------
def ingest_document(openai_client, documents, chunks, spec: dict) -> int:
    print(f"処理開始: {spec['title']}")
    doc_id = _doc_id_for(spec["path"])
    now_iso = datetime.now(timezone.utc).isoformat()
    documents.upsert_item(
        {
            "id": doc_id,
            "title": spec["title"],
            "type": spec["type"],          # パーティションキー
            "industry": spec["industry"],
            "tags": spec["tags"],
            "created_at": now_iso,
            "file_path": spec["path"],
        }
    )
    print(f"  ✓ documents に登録: {doc_id}")

    slides = extract_slides_from_pptx(_DATA_DIR / spec["path"])
    chunk_count = 0
    for slide_number, slide_text in slides:
        for chunk_text in split_into_chunks(slide_text):
            chunks.upsert_item(
                {
                    "id": str(uuid.uuid4()),
                    "document_id": doc_id,         # パーティションキー
                    "text": chunk_text,
                    "embedding": get_embedding(openai_client, chunk_text),
                    "slide_number": slide_number,  # 実スライド番号
                    "section": f"slide_{slide_number}",
                }
            )
            chunk_count += 1
    print(f"  ✓ chunks に登録完了（{chunk_count} 件）\n")
    return chunk_count


def main() -> None:
    openai_client = AzureOpenAI(
        api_key=settings.azure_openai_api_key,
        azure_endpoint=settings.azure_openai_endpoint,
        api_version=settings.azure_openai_api_version,
    )
    db = CosmosClient(
        settings.cosmos_endpoint, credential=settings.cosmos_key
    ).get_database_client(settings.cosmos_database)
    documents, chunks = _setup_containers(db)

    # 取り込み前にクリーンアップ（旧ランダム doc_id データや前回分を一掃して重複防止）。
    # 子（chunks）→ 親（documents）の順で削除する。
    _purge_all(chunks, "document_id", "chunks")
    _purge_all(documents, "type", "documents")

    total = 0
    for spec in SAMPLE_FILES:
        total += ingest_document(openai_client, documents, chunks, spec)
    print(f"[done] ingested {len(SAMPLE_FILES)} documents / {total} chunks")


if __name__ == "__main__":
    main()
