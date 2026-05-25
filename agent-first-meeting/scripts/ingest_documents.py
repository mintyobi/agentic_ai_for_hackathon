"""営業資料 PPTX を documents / chunks コンテナに取り込む（ナレッジベース構築）.

- documents: 資料メタ情報（PK = /type）
- chunks   : 本文チャンク + 1536 次元埋め込み（PK = /document_id, vector policy 付き）

chunks は VectorDistance 検索のためにベクトルポリシーが必須。Cosmos はベクトルポリシーを
コンテナ作成時にしか設定できないため、vector-policy.json / index-policy.json を使って
「無ければベクトルポリシー付きで作成」する。data/ 配下のサンプル PPTX 3 点を投入する。
"""
import json
import sys
import uuid
from datetime import datetime, timezone
from pathlib import Path

from azure.cosmos import CosmosClient, PartitionKey
from openai import AzureOpenAI
from pptx import Presentation

from agent_first_meeting.config import settings

sys.stdout.reconfigure(encoding="utf-8")

_ROOT = Path(__file__).resolve().parent.parent  # agent-first-meeting/
_DATA_DIR = _ROOT / "data"
_VECTOR_POLICY = json.loads((_ROOT / "vector-policy.json").read_text("utf-8"))
_INDEX_POLICY = json.loads((_ROOT / "index-policy.json").read_text("utf-8"))

_CHUNK_SIZE = 500
_CHUNK_OVERLAP = 100

SAMPLE_FILES = [
    {
        "path": "proposal_manufacturing_2024.pptx",
        "title": "製造業向け生産管理DX提案書 2024年版",
        "type": "proposal",
        "industry": "製造業",
        "tags": ["DX", "IoT", "コスト削減", "予防保全"],
    },
    {
        "path": "proposal_it_security_2024.pptx",
        "title": "中堅企業向けセキュリティ強化提案書 2024年版",
        "type": "proposal",
        "industry": "IT",
        "tags": ["セキュリティ", "MFA", "EDR", "ランサムウェア対策"],
    },
    {
        "path": "catalog_cloud_services_2024.pptx",
        "title": "クラウドサービス製品カタログ 2024",
        "type": "catalog",
        "industry": "全業種",
        "tags": ["クラウド", "IoT", "セキュリティ", "データ分析"],
    },
]


def _extract_text_from_pptx(file_path: Path) -> str:
    """スライドごとに本文テキストを抽出する."""
    prs = Presentation(str(file_path))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [
            line
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            if (line := para.text.strip())
        ]
        if lines:
            slides.append(f"[スライド{i}]\n" + "\n".join(lines))
    return "\n\n".join(slides)


def _split(text: str) -> list[str]:
    """500 文字・100 文字オーバーラップの素朴な分割（langchain 非依存）."""
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + _CHUNK_SIZE, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - _CHUNK_OVERLAP
    return chunks


def _setup_containers(db):
    documents = db.create_container_if_not_exists(
        id="documents",
        partition_key=PartitionKey(path="/type"),
    )
    chunks = db.create_container_if_not_exists(
        id="chunks",
        partition_key=PartitionKey(path="/document_id"),
        indexing_policy=_INDEX_POLICY,
        vector_embedding_policy=_VECTOR_POLICY,
    )
    return documents, chunks


def _embed(openai: AzureOpenAI, text: str) -> list[float]:
    return (
        openai.embeddings.create(
            model=settings.azure_openai_embedding_deployment,
            input=text,
            dimensions=1536,
        )
        .data[0]
        .embedding
    )


def ingest_one(openai, documents, chunks, spec: dict) -> int:
    path = _DATA_DIR / spec["path"]
    doc_id = str(uuid.uuid4())
    now_iso = datetime.now(timezone.utc).isoformat()
    documents.upsert_item(
        {
            "id": doc_id,
            "title": spec["title"],
            "type": spec["type"],
            "industry": spec["industry"],
            "tags": spec["tags"],
            "createdAt": now_iso,
            "filePath": spec["path"],
        }
    )
    pieces = _split(_extract_text_from_pptx(path))
    for i, chunk_text in enumerate(pieces):
        chunks.upsert_item(
            {
                "id": str(uuid.uuid4()),
                "document_id": doc_id,
                "text": chunk_text,
                "embedding": _embed(openai, chunk_text),
                "slide_number": i,
                "section": f"chunk_{i}",
            }
        )
    print(f"  ✓ {spec['title']}: doc={doc_id} chunks={len(pieces)}")
    return len(pieces)


def main() -> None:
    openai = AzureOpenAI(
        api_key=settings.azure_openai_api_key,
        azure_endpoint=settings.azure_openai_endpoint,
        api_version=settings.azure_openai_api_version,
    )
    db = CosmosClient(
        settings.cosmos_endpoint, credential=settings.cosmos_key
    ).get_database_client(settings.cosmos_database)
    documents, chunks = _setup_containers(db)

    total = 0
    for spec in SAMPLE_FILES:
        print(f"処理開始: {spec['title']}")
        total += ingest_one(openai, documents, chunks, spec)
    print(f"[done] ingested {len(SAMPLE_FILES)} documents / {total} chunks")


if __name__ == "__main__":
    main()
