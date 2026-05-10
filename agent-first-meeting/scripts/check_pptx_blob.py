"""PowerPoint 表紙生成 + Blob アップロード スモーク (Step 6-7).

python-pptx で表紙 1 枚を生成し、Azure Blob Storage の
generated-documents コンテナにアップロード、ダウンロードで往復確認する。
"""
import sys
from io import BytesIO

from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient
from pptx import Presentation

from agent_first_meeting.config import settings as _settings

sys.stdout.reconfigure(encoding="utf-8")


def make_blob_service_client() -> BlobServiceClient:
    """key があればキー認証、無ければ DefaultAzureCredential を使う."""
    if _settings.blob_account_key:
        return BlobServiceClient(
            account_url=_settings.blob_account_url,
            credential=_settings.blob_account_key,
        )
    return BlobServiceClient(
        account_url=_settings.blob_account_url,
        credential=DefaultAzureCredential(),
    )


COMPANY_NAME = "株式会社サンプル"
TITLE = "DX 推進ご提案"


def build_cover_pptx() -> BytesIO:
    """表紙だけの PowerPoint を BytesIO で返す."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = f"{COMPANY_NAME} 様向け {TITLE}"
    slide.placeholders[1].text = (
        "agent-first-meeting / Microsoft Agent Hackathon 2026"
    )

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def main() -> None:
    print(f"[generate] company={COMPANY_NAME} title={TITLE}")
    pptx_buf = build_cover_pptx()
    pptx_bytes = pptx_buf.getvalue()
    print(f"[pptx] size={len(pptx_bytes)} bytes")

    blob_service = make_blob_service_client()
    auth_mode = "key" if _settings.blob_account_key else "DefaultAzureCredential"
    print(f"[blob_auth] {auth_mode}")
    container_client = blob_service.get_container_client(_settings.blob_container)

    blob_name = "smoke/cover_test.pptx"
    blob_client = container_client.get_blob_client(blob_name)
    blob_client.upload_blob(pptx_bytes, overwrite=True)
    print(f"[upload] OK")
    print(f"[blob_url] {blob_client.url}")

    downloaded = blob_client.download_blob().readall()
    print(f"[download-back] size={len(downloaded)} bytes")
    print(f"[match] {len(downloaded) == len(pptx_bytes)}")


if __name__ == "__main__":
    main()
