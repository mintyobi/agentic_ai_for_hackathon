"""`_pptx_builder.build_presentation_bytes` の構造テスト.

Azure / Semantic Kernel に依存せず、python-pptx だけでローカル実行可能。
"""
from io import BytesIO
from pathlib import Path

from pptx import Presentation

from agent_first_meeting.tools._pptx_builder import (
    AGENDA_ITEMS,
    DEFAULT_PRODUCT_NAME,
    DEFAULT_PRODUCT_PRICE_JPY,
    build_presentation_bytes,
)

# tests/ → agent-first-meeting/ → templates/default.pptx
_DEFAULT_TEMPLATE = Path(__file__).resolve().parents[1] / "templates" / "default.pptx"

# 全テストで使う最小ペイロード
SAMPLE_KWARGS = dict(
    cover_title="製造業のDX：技能継承課題への AI ナレッジ活用ご提案",
    cover_subtitle="株式会社サンプル製作所 様向け / 2026年5月 / 担当: 佐々木",
    industry_body="ベテラン技術者の高齢化\n暗黙知の継承困難\n人手不足",
    position_body="経営層: 人材リスクの定量化\n部門責任者: KPI 影響\n担当者: 現場 UI",
)


def _open_generated() -> Presentation:
    data = build_presentation_bytes(**SAMPLE_KWARGS)
    assert isinstance(data, bytes)
    assert len(data) > 5_000, "PPTX が小さすぎる（破損？）"
    return Presentation(BytesIO(data))


def test_six_slides_are_generated():
    """スライド枚数が固定の 6 枚であること."""
    prs = _open_generated()
    assert len(prs.slides) == 6


def test_slide_titles_are_in_expected_order():
    """各スライドのタイトル順が想定どおりであること."""
    prs = _open_generated()
    titles = [slide.shapes.title.text for slide in prs.slides]
    assert titles == [
        SAMPLE_KWARGS["cover_title"],
        "本日のアジェンダ",
        "業界トレンドとお客様の課題",
        "ご担当者向けのご提案",
        "自社商品のご紹介",
        "費用について",
    ]


def test_cover_subtitle_is_rendered():
    """表紙のサブタイトルが配置されていること."""
    prs = _open_generated()
    cover = prs.slides[0]
    # placeholder[1] = サブタイトル
    assert cover.placeholders[1].text == SAMPLE_KWARGS["cover_subtitle"]


def test_industry_body_is_split_into_bullets():
    """業界スライドの本文が改行で箇条書きに分割されていること."""
    prs = _open_generated()
    industry_slide = prs.slides[2]
    body_text = industry_slide.placeholders[1].text
    assert "ベテラン技術者の高齢化" in body_text
    assert "暗黙知の継承困難" in body_text
    assert "人手不足" in body_text
    # 改行が保持されている（= 箇条書きとして展開されている）
    lines = body_text.splitlines()
    assert len(lines) >= 3


def test_position_body_renders_role_specific_points():
    """役職向けスライドが取引相手の役職別論点を含むこと."""
    prs = _open_generated()
    position_slide = prs.slides[3]
    body_text = position_slide.placeholders[1].text
    assert "経営層" in body_text
    assert "部門責任者" in body_text
    assert "担当者" in body_text


def test_product_slide_uses_default_product_name():
    """自社商品スライドに既定の商品名（プレースホルダ）が入ること."""
    prs = _open_generated()
    product_slide = prs.slides[4]
    body_text = product_slide.placeholders[1].text
    assert DEFAULT_PRODUCT_NAME in body_text
    # 偽の商品名（テスト商品）を埋め込まないこと
    assert DEFAULT_PRODUCT_NAME != "テスト商品"


def test_cost_slide_omits_fake_price_when_unset():
    """既定価格 0 のときは偽の金額を出さず「別途お見積もり」にすること."""
    prs = _open_generated()
    cost_slide = prs.slides[5]
    body_text = cost_slide.placeholders[1].text
    assert DEFAULT_PRODUCT_PRICE_JPY == 0
    assert "別途お見積もり" in body_text
    assert "10 円" not in body_text
    assert DEFAULT_PRODUCT_NAME in body_text


def test_product_price_can_be_overridden():
    """既定値を上書きできること（将来の本実装に向けた契約テスト）."""
    data = build_presentation_bytes(
        **SAMPLE_KWARGS,
        product_name="プレミアム商品",
        product_price_jpy=120_000,
    )
    prs = Presentation(BytesIO(data))
    cost_body = prs.slides[5].placeholders[1].text
    product_body = prs.slides[4].placeholders[1].text
    assert "プレミアム商品" in cost_body
    assert "プレミアム商品" in product_body
    # 桁区切りで表示される
    assert "120,000" in cost_body


def test_agenda_items_are_rendered_in_table_of_contents():
    """目次スライドが AGENDA_ITEMS をすべて含むこと."""
    prs = _open_generated()
    toc_slide = prs.slides[1]
    toc_text = toc_slide.placeholders[1].text
    for item in AGENDA_ITEMS:
        assert item in toc_text


# ---------- テンプレ（自社フォーマット）モード ----------

def _open_with_template() -> Presentation:
    """templates/default.pptx を使って生成した PPTX を返す."""
    data = build_presentation_bytes(
        **SAMPLE_KWARGS, template_path=str(_DEFAULT_TEMPLATE)
    )
    return Presentation(BytesIO(data))


def test_template_mode_produces_six_slides():
    """テンプレ経路でも 6 スライド生成されること（テンプレ内見本スライドは除去される）."""
    assert _DEFAULT_TEMPLATE.exists(), f"{_DEFAULT_TEMPLATE} が存在しない"
    prs = _open_with_template()
    assert len(prs.slides) == 6


def test_template_mode_uses_named_layouts():
    """default.pptx では表紙が "Cover"（命名）、本文は fallback の "Title and Content" を使うこと.

    自社の独自デザインを当てたい場合は、対象レイアウトを PowerPoint で複製し
    Agenda / Industry / Position / Product / Cost に改名すれば本コードが自動で拾う。
    """
    prs = _open_with_template()
    actual_names = [slide.slide_layout.name for slide in prs.slides]
    assert actual_names[0] == "Cover"
    # 本文5枚は標準名 "Title and Content" にフォールバックして同一レイアウトを共有
    assert actual_names[1:] == ["Title and Content"] * 5


def test_template_mode_content_matches_no_template():
    """テンプレ経路でも、表紙タイトル・本文・商品名などの内容が同等に流し込まれること."""
    prs = _open_with_template()
    assert prs.slides[0].shapes.title.text == SAMPLE_KWARGS["cover_title"]
    assert prs.slides[0].placeholders[1].text == SAMPLE_KWARGS["cover_subtitle"]
    industry_body = prs.slides[2].placeholders[1].text
    assert "ベテラン技術者の高齢化" in industry_body
    cost_body = prs.slides[5].placeholders[1].text
    assert DEFAULT_PRODUCT_NAME in cost_body
