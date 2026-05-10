"""PowerPoint 生成プラグイン (python-pptx + Azure Blob)."""
import uuid
from datetime import datetime, timedelta, timezone
from io import BytesIO
from typing import Annotated

from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobClient, BlobSasPermissions, BlobServiceClient, generate_blob_sas
from pptx import Presentation
from semantic_kernel.functions import kernel_function

from agent_first_meeting.config import settings


SAS_EXPIRY_HOURS = 24


def _make_blob_service_client() -> BlobServiceClient:
    if settings.blob_account_key:
        return BlobServiceClient(
            account_url=settings.blob_account_url,
            credential=settings.blob_account_key,
        )
    return BlobServiceClient(
        account_url=settings.blob_account_url,
        credential=DefaultAzureCredential(),
    )


class DocumentGenPlugin:
    """表紙 1 枚の PowerPoint を生成し Blob にアップロードする SK プラグイン."""

    def __init__(self) -> None:
        self._blob_service = _make_blob_service_client()
        self._container = self._blob_service.get_container_client(
            settings.blob_container
        )

    def _build_download_url(self, blob_client: BlobClient) -> str:
        """アップロードした Blob のダウンロード URL を返す.

        - blob_account_key がある場合：account-key SAS を付与した時限 URL
        - 無い場合（Managed Identity 想定）：素の URL を返す
          ※ 本番では get_user_delegation_key() ベースの SAS に置換予定
        """
        if not settings.blob_account_key:
            return blob_client.url

        sas_token = generate_blob_sas(
            account_name=self._blob_service.account_name,
            container_name=settings.blob_container,
            blob_name=blob_client.blob_name,
            account_key=settings.blob_account_key,
            permission=BlobSasPermissions(read=True),
            expiry=datetime.now(timezone.utc) + timedelta(hours=SAS_EXPIRY_HOURS),
        )
        return f"{blob_client.url}?{sas_token}"

    @kernel_function(
        description=(
            "表紙だけの初回提案資料 (PowerPoint) を生成し、"
            "Azure Blob にアップロードしてダウンロード可能な URL を返す。"
        ),
    )
    def generate_pptx(
        self,
        title: Annotated[
            str,
            "表紙のメインタイトル。例: '製造業のDX：技能継承課題への AI ナレッジ活用ご提案'",
        ],
        subtitle: Annotated[
            str,
            "表紙のサブタイトル。例: '株式会社サンプル製作所 様向け / 2026年5月'",
        ] = "",
    ) -> Annotated[str, "生成された PowerPoint の Blob URL."]:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

        buf = BytesIO()
        prs.save(buf)

        blob_name = f"proposals/{uuid.uuid4().hex}.pptx"
        blob_client = self._container.get_blob_client(blob_name)
        blob_client.upload_blob(buf.getvalue(), overwrite=True)
        return self._build_download_url(blob_client)
