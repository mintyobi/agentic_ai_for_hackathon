"""PPTX 組み立てロジック（Azure 非依存・ユニットテスト容易）.

`document_gen.py` の Blob アップロード処理から分離し、
ローカルだけで PPTX 生成が再現できるようにする。

自社フォーマットのテンプレート（`.pptx`）を指定すると、その**スライドマスター・
レイアウト・配色・フォント**を引き継いで描画する。テンプレ無しなら python-pptx の
既定 `Presentation()` で描画（従来挙動）。

レイアウトの選択ルール（テンプレ作成側との契約）:
  - 表紙: 名前 "Cover" を探し、無ければ python-pptx 既定の "Title Slide" にフォールバック
  - 本文6スライド: 名前 "Agenda" / "Industry" / "Position" / "Product" / "Cost"
    を順に探し、無ければ "Title and Content" にフォールバック
  - これらも無ければそれぞれインデックス 0 / 1 にフォールバック
テンプレ側は PowerPoint で自由にデザインでき、レイアウトに上記の名前を付けるか
あるいは標準名を保てば、本コードが正しく流し込みます。
"""
from io import BytesIO

from pptx import Presentation
from pptx.presentation import Presentation as _Presentation  # 型注釈用

# 「自社商品 / 費用」スライドの既定文言。実値は config（.env）で上書きする想定。
# 価格 0 以下は「金額未設定」とみなし、偽の金額の代わりに「別途お見積もり」を出す。
DEFAULT_PRODUCT_NAME = "弊社ソリューション"
DEFAULT_PRODUCT_PRICE_JPY = 0
AGENDA_ITEMS = [
    "1. ご挨拶",
    "2. 業界トレンドとお客様の課題",
    "3. ご担当者向けのご提案",
    "4. 自社商品のご紹介",
    "5. 費用について",
]

# 各スライドが使うレイアウトの解決ルール（名前候補, 既定インデックス）
_LAYOUT_RESOLUTION = {
    "cover":    (["Cover", "Title Slide"], 0),
    "agenda":   (["Agenda", "Title and Content"], 1),
    "industry": (["Industry", "Title and Content"], 1),
    "position": (["Position", "Title and Content"], 1),
    "product":  (["Product", "Title and Content"], 1),
    "cost":     (["Cost", "Title and Content"], 1),
}


def _resolve_layout(prs: _Presentation, key: str):
    """名前候補→インデックスの順でレイアウトを解決."""
    names, fallback_idx = _LAYOUT_RESOLUTION[key]
    by_name = {layout.name: layout for layout in prs.slide_layouts}
    for name in names:
        if name in by_name:
            return by_name[name]
    return prs.slide_layouts[fallback_idx]


def _clear_existing_slides(prs: _Presentation) -> None:
    """テンプレに残っている見本スライドを全て除去（マスター・レイアウトは温存）.

    python-pptx には公開APIが無いので、内部の `_sldIdLst` から sldId 要素を外す
    定石パターンを使う。/ppt/slides/* の中身は残るが参照されないため実害なし。
    """
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    for sld_id in list(sld_id_lst):
        sld_id_lst.remove(sld_id)


def _add_cover(prs: _Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(_resolve_layout(prs, "cover"))
    slide.shapes.title.text = title
    # サブタイトル placeholder（Title Slide の標準は idx=1）
    slide.placeholders[1].text = subtitle


def _add_content_slide(prs: _Presentation, key: str, title: str, body: str) -> None:
    """Title + bullet body の標準スライドを追加."""
    slide = prs.slides.add_slide(_resolve_layout(prs, key))
    slide.shapes.title.text = title
    # placeholder[1] は body プレースホルダ。複数行は自動で箇条書きになる
    slide.placeholders[1].text = body


def build_presentation_bytes(
    cover_title: str,
    cover_subtitle: str,
    industry_body: str,
    position_body: str,
    product_name: str = DEFAULT_PRODUCT_NAME,
    product_price_jpy: int = DEFAULT_PRODUCT_PRICE_JPY,
    template_path: str | None = None,
) -> bytes:
    """6 スライド構成の PPTX をメモリ上に生成して bytes で返す.

    Azure に依存しない純粋関数。ローカルでユニットテスト可能。
    `template_path` を渡すと、その `.pptx` のマスター/レイアウト/配色を引き継ぐ。
    """
    prs = Presentation(template_path) if template_path else Presentation()
    if template_path:
        # テンプレに残った見本スライドを除去（マスター/レイアウトは保持）
        _clear_existing_slides(prs)

    # 1. 表紙
    _add_cover(prs, cover_title, cover_subtitle)

    # 2. 目次
    _add_content_slide(prs, "agenda", "本日のアジェンダ", "\n".join(AGENDA_ITEMS))

    # 3. 対業界向け
    _add_content_slide(prs, "industry", "業界トレンドとお客様の課題", industry_body)

    # 4. 役職向け
    _add_content_slide(prs, "position", "ご担当者向けのご提案", position_body)

    # 5. 自社商品
    _add_content_slide(
        prs,
        "product",
        "自社商品のご紹介",
        f"商品名：{product_name}\n本日はこちらの商品をご提案いたします。",
    )

    # 6. 費用（価格未設定なら偽の金額を出さず「別途お見積もり」とする）
    price_line = (
        f"価格：{product_price_jpy:,} 円（税別）"
        if product_price_jpy > 0
        else "価格：別途お見積もり"
    )
    _add_content_slide(
        prs,
        "cost",
        "費用について",
        (
            f"商品「{product_name}」のご提供価格\n"
            f"{price_line}\n"
            "※ 詳細条件は別途ご相談のうえ決定いたします。"
        ),
    )

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
